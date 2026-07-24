import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation_12():
    prs = Presentation()
    
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Color Palette (High-Contrast Light Theme)
    BG_COLOR = RGBColor(248, 250, 252)       # Light Slate #f8fafc
    CARD_BG = RGBColor(255, 255, 255)        # Pure White #ffffff
    CARD_BORDER = RGBColor(203, 213, 225)    # Border Slate #cbd5e1
    TEXT_DARK = RGBColor(15, 23, 42)         # Primary Text #0f172a
    TEXT_MUTED = RGBColor(71, 85, 105)       # Muted Text #475569
    ACCENT_BLUE = RGBColor(37, 99, 235)      # Primary Blue #2563eb
    ACCENT_GREEN = RGBColor(5, 150, 105)     # Emerald Green #059669
    ACCENT_DANGER = RGBColor(220, 38, 38)     # Warning Red #dc2626
    ACCENT_AMBER = RGBColor(217, 119, 6)      # Amber Yellow #d97706
    
    FONT_TITLE = 'Noto Sans JP'
    FONT_BODY = 'Noto Sans JP'

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title_text, subtitle_text=None):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.833), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_TITLE
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = FONT_BODY
            p2.font.size = Pt(15)
            p2.font.color.rgb = TEXT_MUTED
            p2.space_before = Pt(4)

    def add_bullets(tf, bullets):
        for item in bullets:
            p_item = tf.add_paragraph()
            p_item.text = item
            p_item.font.name = FONT_BODY
            p_item.font.size = Pt(17)
            p_item.font.color.rgb = TEXT_DARK
            p_item.space_after = Pt(12)

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "現場改善プロジェクト報告"
    p0.font.name = FONT_BODY
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_BLUE
    p0.space_after = Pt(10)
    
    p1 = tf.add_paragraph()
    p1.text = "Seibi"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_BLUE
    p1.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "点検のデジタル化からAIによる修理サポートまで"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_DARK
    p2.space_after = Pt(30)
    
    p3 = tf.add_paragraph()
    p3.text = "溶接ラインチーム"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 2: 背景と課題
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_header(slide, "背景と課題", "紙の点検表運用による3つの現場課題")
    
    card_w = Inches(3.64)
    card_h = Inches(4.5)
    
    problems = [
        ("📁 履歴の散逸", "手書きの点検表は提出後に保管または破棄され、検索可能なデジタル記録が残らない状態でした。"),
        ("🔄 追跡の困難", "同じトラブルが繰り返し発生しても記憶に頼るしかなく、発見された異常が確実に修理されたか確認できませんでした。"),
        ("📖 作業の非効率", "トラブル対応時に紙のマニュアルを探してめくるか、特定の人に質問するしか選択肢がありませんでした。")
    ]
    
    for idx, (p_title, p_desc) in enumerate(problems):
        left_pos = Inches(0.75 + idx * 3.94)
        c = slide.shapes.add_textbox(left_pos, Inches(1.8), card_w, card_h)
        tf_c = c.text_frame
        tf_c.word_wrap = True
        
        p = tf_c.paragraphs[0]
        p.text = p_title
        p.font.name = FONT_TITLE
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.space_after = Pt(14)
        
        p_sub = tf_c.add_paragraph()
        p_sub.text = p_desc
        p_sub.font.name = FONT_BODY
        p_sub.font.size = Pt(16)
        p_sub.font.color.rgb = TEXT_DARK
        p_sub.line_spacing = 1.4

    # ==========================================
    # SLIDE 3: 解決策と取り組み (Solutions Intro)
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_header(slide, "解決策と取り組み", "4つの主要機能で現場の課題を解決")
    
    agenda_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.0), Inches(8.333), Inches(4.5))
    tf_agenda = agenda_box.text_frame
    tf_agenda.word_wrap = True
    
    solutions_agenda = [
        "1. 点検業務のデジタル化と履歴保存",
        "2. レイアウトマップによる視覚管理",
        "3. リアルタイムの状況共有とタスク管理",
        "4. AI搭載の修理サポート機能"
    ]
    
    for idx, item in enumerate(solutions_agenda):
        p_item = tf_agenda.paragraphs[0] if idx == 0 else tf_agenda.add_paragraph()
        p_item.text = item
        p_item.font.name = FONT_BODY
        p_item.font.size = Pt(22)
        p_item.font.bold = True
        p_item.font.color.rgb = TEXT_DARK
        p_item.space_after = Pt(24)

    # Helper for split layout
    def add_split_slide(title, subtitle, bullets, img_name):
        s = prs.slides.add_slide(blank_layout)
        set_background(s)
        add_header(s, title, subtitle)
        
        left_box = s.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.8), Inches(5.0))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        add_bullets(tf_left, bullets)
        
        if os.path.exists(img_name):
            s.shapes.add_picture(img_name, Inches(6.8), Inches(1.8), width=Inches(5.7))

    # ==========================================
    # SLIDE 4: 解決策(1) 点検業務のデジタル化
    # ==========================================
    add_split_slide(
        "解決策：点検業務のデジタル化",
        "モバイル端末でその場入力・ペーパーレス化",
        [
            "📱 現場でのモバイル入力: スマホやタブレットから直接点検表を入力。転記ミスや紛失を防止します。",
            "☁️ 検索可能な履歴の蓄積: 点検データや過去の不具合履歴がクラウドへ保存され、いつでも検索可能です。"
        ],
        "app_dashboard.png"
    )

    # ==========================================
    # SLIDE 5: 解決策：検索可能な履歴ログ
    # ==========================================
    add_split_slide(
        "解決策：検索可能な履歴ログ",
        "過去の不具合や修理記録を瞬時に検索",
        [
            "🔍 即座のキーワード検索: 日付や設備名、症状を入力するだけで、過去の修理履歴をその場で参照できます。",
            "📊 再発トラブルの把握: 過去にどんな対応をしたかが可視化され、熟練者に頼らず迅速な対応を可能にします。"
        ],
        "app_history.png"
    )

    # ==========================================
    # SLIDE 6: 解決策：レイアウトマップによる視覚管理 [RESTORED]
    # ==========================================
    add_split_slide(
        "解決策：レイアウトマップによる視覚管理",
        "平面図上で設備配置と配線ルートを直感的に把握",
        [
            "🗺️ インタラクティブマップ: ロボットやレギュレーター等の配置や異常ステータスを、マップ上で一目に視認可能。",
            "🔌 配線ルートの可視化: 機器をタップするだけで配線系統を表示し、トラブル箇所の特定をスピード化。"
        ],
        "app_wire_map.png"
    )

    # ==========================================
    # SLIDE 7: 解決策：状況共有と管理
    # ==========================================
    add_split_slide(
        "解決策：状況共有と管理",
        "自動タスク登録とLINE WORKS通知で連絡漏れを防止",
        [
            "🔔 自動タスク化＆通知: 異常「あり」と回答されると自動で修理タスクが登録され、LINE WORKSへ即座に緊急通知されます。",
            "🔄 全デバイス同期: 修理の進捗状況や掲示板の内容が、すべてのデバイス（モバイル・PC）へリアルタイムに同期されます。"
        ],
        "app_bulletin.png"
    )

    # ==========================================
    # SLIDE 8: 解決策：AI修理サポート
    # ==========================================
    add_split_slide(
        "解決策：AI修理サポート",
        "現場ですぐに頼れるAIアドバイザー機能",
        [
            "🤖 普段の言葉で質問応答: マニュアルをめくる代わりに、自然な言葉でトラブル症状を質問できます。",
            "📄 明確な引用元ページ提示: AIは取り込んだマニュアルを参照し、根拠となるページや項番を示して回答します。"
        ],
        "app_manuals.png"
    )

    # ==========================================
    # SLIDE 9: 成果と今後の展望 (Outlook Intro)
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_header(slide, "成果と今後の展望", "さらなる現場改善に向けた3つのステップ")
    
    agenda_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.0), Inches(8.333), Inches(4.5))
    tf_agenda = agenda_box.text_frame
    tf_agenda.word_wrap = True
    
    outlook_agenda = [
        "1. データ主導の予防保全",
        "2. AI対応マニュアルの拡大",
        "3. 使うほど賢くなるAI学習ループ"
    ]
    
    for idx, item in enumerate(outlook_agenda):
        p_item = tf_agenda.paragraphs[0] if idx == 0 else tf_agenda.add_paragraph()
        p_item.text = item
        p_item.font.name = FONT_BODY
        p_item.font.size = Pt(22)
        p_item.font.bold = True
        p_item.font.color.rgb = TEXT_DARK
        p_item.space_after = Pt(24)

    # ==========================================
    # SLIDE 10: 今後の展望：データ主導の保全
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_header(slide, "今後の展望：データ主導の保全", "点検ログの傾向から故障を予測し未然防止")
    
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.8), Inches(5.0))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    add_bullets(tf_left, [
        "⚡ 溶接機ログの予兆検知: 点検表の「ワイヤ送給のもたつき」や「滑り」の頻度から、ライン停止のサインを検知。",
        "🛠️ 予測交換によるダウンタイムゼロ: ライナーの詰まりや駆動部の摩耗期を特定し、次回段取り時に交換してトラブルを防止。"
    ])
    
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    p_r0 = tf_right.paragraphs[0]
    p_r0.text = "溶接機 状態モニター（過去5ヶ月間）"
    p_r0.font.name = FONT_TITLE
    p_r0.font.size = Pt(20)
    p_r0.font.bold = True
    p_r0.font.color.rgb = ACCENT_BLUE
    p_r0.space_after = Pt(16)
    
    sim_data = [
        "🔴 ワイヤ送給のもたつき: 3回検出 [警告基準: 2回以上]",
        "🟡 ワイヤの滑り: 2回検出 [注意基準: 1回以上]",
        "🔧 推奨アクション: 点検記録より「もたつき」「滑り」の再発傾向を検知。金属粉による詰まりを防ぐため、次回段取り替え時にコンジットライナー交換を推奨します。"
    ]
    add_bullets(tf_right, sim_data)

    # ==========================================
    # SLIDE 11: 今後の展望：AI機能の高度化
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_header(slide, "今後の展望：AI機能の高度化", "全設備のカバーとAI自律学習の確立")
    
    box_w = Inches(5.7)
    box_h = Inches(2.2)
    
    c1 = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), box_w, box_h)
    tf1 = c1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "📚 AI対応マニュアルの拡大"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_BLUE
    p1.space_after = Pt(10)
    add_bullets(tf1, ["工場内のすべての主要設備の手順書・図面データをAIへ順次取り込み、全域をサポート可能にします。"])
    
    c2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), box_w, box_h)
    tf2 = c2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "🔄 使うほど賢くなるAI学習ループ"
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN
    p2.space_after = Pt(10)
    add_bullets(tf2, ["質問ログとフィードバック（役に立ったか）を分析。良回答を保存・学習し、精度を向上させます。"])
    
    c3 = slide.shapes.add_textbox(Inches(0.75), Inches(4.5), Inches(11.833), Inches(1.8))
    tf3 = c3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "💡 まとめ"
    p3.font.name = FONT_TITLE
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = TEXT_DARK
    p3.space_after = Pt(8)
    add_bullets(tf3, ["Seibi: 製造現場のダウンタイムをゼロにし、組織の保全力と稼働率を最大化するソリューション"])

    # ==========================================
    # SLIDE 12: 閉会
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    
    closing_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(3.0))
    tf_closing = closing_box.text_frame
    tf_closing.word_wrap = True
    
    p_close = tf_closing.paragraphs[0]
    p_close.text = "ご清聴ありがとうございました"
    p_close.font.name = FONT_TITLE
    p_close.font.size = Pt(52)
    p_close.font.bold = True
    p_close.font.color.rgb = ACCENT_GREEN
    p_close.space_after = Pt(16)
    
    p_sub_close = tf_closing.add_paragraph()
    p_sub_close.text = "設備管理システム — Seibi"
    p_sub_close.font.name = FONT_BODY
    p_sub_close.font.size = Pt(24)
    p_sub_close.font.color.rgb = TEXT_MUTED

    # Save presentation
    output_path = 'seibi_presentation.pptx'
    prs.save(output_path)
    print(f"12-Slide PowerPoint created successfully at: {output_path}")

if __name__ == '__main__':
    create_presentation_12()
