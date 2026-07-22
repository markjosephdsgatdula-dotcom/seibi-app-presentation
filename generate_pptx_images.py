import os
from pptx import Presentation
from pptx.util import Inches

def create_visual_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]  # Blank slide layout
    
    for i in range(1, 12):
        img_path = f"web_slide_{i}.png"
        if os.path.exists(img_path):
            slide = prs.slides.add_slide(blank_layout)
            # Add image stretching across the entire slide (full-bleed)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            print(f"Added Slide {i} using the high-resolution HTML screenshot.")
        else:
            print(f"Warning: {img_path} not found.")
            
    output_path = 'seibi_presentation_visual.pptx'
    prs.save(output_path)
    print(f"Visual presentation successfully compiled at: {output_path}")

if __name__ == '__main__':
    create_visual_presentation()
