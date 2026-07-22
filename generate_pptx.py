import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Design constants (Seibi App Theme)
    BG_COLOR = RGBColor(11, 14, 20)       # Dark background
    TEXT_LIGHT = RGBColor(243, 244, 246)  # Off-white
    TEXT_MUTED = RGBColor(156, 163, 175)  # Muted grey
    ACCENT_BLUE = RGBColor(59, 130, 246)   # Bright blue
    ACCENT_GREEN = RGBColor(16, 185, 129) # Success green
    ACCENT_RED = RGBColor(239, 68, 68)    # Alert red
    
    # Using standard Japanese fonts available on Windows
    FONT_TITLE = "Yu Gothic"
    FONT_BODY = "Yu Gothic"
    
    # Helper to set slide background
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    # Helper to create a slide title
    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_TITLE
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
    # Helper to format bullets nicely
    def add_bullets(text_frame, items, is_muted=False):
        for idx, item in enumerate(items):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = item
            p.font.name = FONT_BODY
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_MUTED if is_muted else TEXT_LIGHT
            p.space_after = Pt(10)
            p.level = 0

    blank_layout = prs.slide_layouts[6]
    
    # ==========================================
    # SLIDE 1: タイトルスライド (Title Slide)
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    
    # Large Logo Title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Seibi"
    p.font.name = FONT_TITLE
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    p2 = tf.add_paragraph()
    p2.text = "設備点検・不具合管理アプリのご紹介"
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(24)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(10)
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.333), Inches(1.5))
    tf_sub = sub_box.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "点検業務のデジタル化から、リアルタイム状況共有、AI修理サポートまで"
    p_sub.font.name = FONT_BODY
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = TEXT_MUTED
    
    p_sub2 = tf_sub.add_paragraph()
    p_sub2.text = "溶接チーム マーク"
    p_sub2.font.name = FONT_BODY
    p_sub2.font.size = Pt(14)
    p_sub2.font.color.rgb = ACCENT_GREEN
    p_sub2.space_before = Pt(5)

    # ==========================================
    # SLIDE 2: 背景と課題 (The Background)
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_slide_header(slide, "背景と課題")
    
    # Left Content - The Problem
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    p = tf_left.paragraphs[0]
    p.text = "従来の紙ベース運用の限界"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED
    p.space_after = Pt(20)
    
    problems = [
        "❌ 点検履歴の散逸: 手書きの点検表は提出後にファイル保管されるか処分され、後から検索できるデジタル記録が残らない。",
        "❌ 過去の状況確認不可: 過去にどんな不具合があったのか、先月や昨年に設備に何が起きたのかを確認する方法がない。"
    ]
    add_bullets(tf_left, problems)

    # Right Content - Operational Blocks
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    p_right = tf_right.paragraphs[0]
    p_right.text = "現場で発生していた具体的な問題"
    p_right.font.name = FONT_TITLE
    p_right.font.size = Pt(20)
    p_right.font.bold = True
    p_right.font.color.rgb = TEXT_LIGHT
    p_right.space_after = Pt(20)
    
    opportunities = [
        "⚠️ 繰り返す不具合の放置: 不具合が再発しても個人の記憶に頼るしかなく、不具合が実際に修理されたかを追跡する仕組みもない。",
        "⚠️ マニュアル参照の非効率: 修理中に手順を確認したい場合、分厚い紙マニュアルをめくるか周りに聞くしかない。"
    ]
    add_bullets(tf_right, opportunities)

    # ==========================================
    # SLIDE 3: 解決策と取り組み (Separator Agenda Slide)
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    
    # Large Section Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "解決策と取り組み"
    p.font.name = FONT_TITLE
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    # Topic List
    list_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.2), Inches(11.833), Inches(3.0))
    tf_list = list_box.text_frame
    tf_list.word_wrap = True
    
    agenda_solutions = [
        "1. 点検業務のデジタル化",
        "2. リアルタイムの状況共有と管理",
        "3. AI搭載の修理サポート機能"
    ]
    for idx, item in enumerate(agenda_solutions):
        if idx == 0:
            p_item = tf_list.paragraphs[0]
        else:
            p_item = tf_list.add_paragraph()
        p_item.text = item
        p_item.font.name = FONT_BODY
        p_item.font.size = Pt(22)
        p_item.font.bold = True
        p_item.font.color.rgb = TEXT_LIGHT
        p_item.space_after = Pt(20)

    # ==========================================
    # SLIDE 4: 解決策：点検業務のデジタル化 (Digitalization)
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_slide_header(slide, "解決策：点検業務のデジタル化")
    
    desc_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(4.5), Inches(4.5))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p = tf_desc.paragraphs[0]
    p.text = "現場でのモバイル入力"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(15)
    
    bullets = [
        "📱 タブレットやスマホでの入力: 現場作業を行いながらその場から直接アプリに入力し提出可能。",
        "🔒 データの安全な保存: 転記ミスや紛失を解消し、すべての点検結果は即座にデジタルデータとして安全に保存。"
    ]
    add_bullets(tf_desc, bullets)
    
    img_path = 'app_dashboard.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(5.6), Inches(1.6), width=Inches(7.0))
        
    # ==========================================
    # SLIDE 5: 解決策：検索可能な履歴ログ (History)
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_slide_header(slide, "解決策：検索可能な履歴ログ")
    
    desc_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(4.5), Inches(4.5))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p = tf_desc.paragraphs[0]
    p.text = "設備ログのデータベース化"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(15)
    
    bullets = [
        "🔍 過去ログの動的検索: 過去の不具合履歴や点検結果を、日付や設備名から瞬時に呼び出せる検索機能を搭載。",
        "📈 不具合履歴の可視化: 設備ごとに不具合発生パターンを蓄積し、同じトラブルの再発傾向をデジタル上で一目で把握。"
    ]
    add_bullets(tf_desc, bullets)
    
    img_path = 'app_history.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(5.6), Inches(1.6), width=Inches(7.0))

    # ==========================================
    # SLIDE 6: 解決策：状況共有と管理 (Notice Board)
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_slide_header(slide, "解決策：状況共有と管理")
    
    desc_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(4.5), Inches(4.5))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p = tf_desc.paragraphs[0]
    p.text = "連絡漏れを防ぐ自動連携"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(15)
    
    bullets = [
        "⚙️ 自動タスク化: 点検表で「異常あり」と回答されると、システムが自動で修理タスクを作成し即座に登録。",
        "🚨 LINE WORKSへの自動アラート: 重大な不具合が検知された場合、LINE WORKSへ自動で緊急通知を送信。",
        "🔄 リアルタイム同期: 修理の状況や掲示板の内容が、すべてのデバイス（モバイル・PC）へリアルタイムに同期。"
    ]
    add_bullets(tf_desc, bullets)
    
    img_path = 'app_bulletin.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(5.6), Inches(1.6), width=Inches(7.0))

    # ==========================================
    # SLIDE 7: 解決策：AI修理サポート (AI Manuals)
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_slide_header(slide, "解決策：AI修理サポート")
    
    desc_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(4.5), Inches(4.5))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p = tf_desc.paragraphs[0]
    p.text = "現場の相談役としてのAI"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(15)
    
    bullets = [
        "💬 自然言語でのスマート質問: 紙マニュアルをめくる代わりに、普段の言葉で質問可能（例：「ワイヤ送給がもたついている」）。",
        "📖 引用元の提示: AIがマニュアルから関連箇所を読み取り、参照したページや項番を示しながら明確な回答。",
        "⚡ トラブル一次対応の迅速化: 直接先輩社員に聞いたり、電話をする手間を省き、その場で即時に一次対処が可能。"
    ]
    add_bullets(tf_desc, bullets)
    
    img_path = 'app_manuals.png'
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(5.6), Inches(1.6), width=Inches(7.0))

    # ==========================================
    # SLIDE 8: 成果と今後の展望 (Separator Agenda Slide)
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    
    # Large Section Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.833), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "成果と今後の展望"
    p.font.name = FONT_TITLE
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    # Topic List
    list_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.2), Inches(11.833), Inches(3.0))
    tf_list = list_box.text_frame
    tf_list.word_wrap = True
    
    agenda_outlook = [
        "1. データ主導のメンテナンス",
        "2. AI対応マニュアルの拡大",
        "3. 使うほど賢くなるAI学習ループ"
    ]
    for idx, item in enumerate(agenda_outlook):
        if idx == 0:
            p_item = tf_list.paragraphs[0]
        else:
            p_item = tf_list.add_paragraph()
        p_item.text = item
        p_item.font.name = FONT_BODY
        p_item.font.size = Pt(22)
        p_item.font.bold = True
        p_item.font.color.rgb = TEXT_LIGHT
        p_item.space_after = Pt(20)

    # ==========================================
    # SLIDE 9: 今後の展望：データ主導の保全 (Predictive Maintenance)
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_slide_header(slide, "今後の展望：データ主導の保全")
    
    # Left Content - Welding Case Study
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    p = tf_left.paragraphs[0]
    p.text = "溶接機の予兆検知とスケジュール交換"
    p.font.name = FONT_TITLE
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.space_after = Pt(15)
    
    welding_bullets = [
        "📊 長期トレンドの統計分析: 蓄積された履歴を基に長期的な故障サイクルを分析し、最適なスケジュールメンテナンスを実現。",
        "🔩 溶接機コンジットライナー詰まり予測: 点検記録からワイヤ送給のもたつきや滑り、アーク切れの予兆を特定。",
        "⚙️ ドライブロールの摩耗管理: 金属粉の蓄積やロール摩耗時期をデータで特定し、最適なタイミングで交換を実施してアーク停止を防止。"
    ]
    add_bullets(tf_left, welding_bullets)

    # Right Content - Simulated dashboard card
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    p_right = tf_right.paragraphs[0]
    p_right.text = "点検履歴分析（過去5ヶ月間）"
    p_right.font.name = FONT_TITLE
    p_right.font.size = Pt(20)
    p_right.font.bold = True
    p_right.font.color.rgb = ACCENT_BLUE
    p_right.space_after = Pt(20)
    
    sim_data = [
        "🔴 ワイヤ送給のもたつき: 3回検出 [警告基準: 2回以上]",
        "🟡 ワイヤの滑り: 2回検出 [注意基準: 1回以上]",
        "🔧 推奨アクション: 点検記録より「もたつき」「滑り」の再発傾向を検知。金属粉による詰まりを防ぐため、次回の段取り替え時にライナー交換を推奨します。"
    ]
    add_bullets(tf_right, sim_data)

    # ==========================================
    # SLIDE 10: 今後の展望：AI機能の高度化
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide)
    add_slide_header(slide, "今後の展望：AI機能の高度化")
    
    box_w = Inches(5.5)
    box_h = Inches(2.2)
    
    # Card 1: AI manual expansion
    c1 = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), box_w, box_h)
    tf1 = c1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "📚 AI対応マニュアルの拡大"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_BLUE
    p1.space_after = Pt(10)
    add_bullets(tf1, [
        "一部の設備だけでなく、工場内のすべての主要設備の手順書・図面データをAIへ順次取り込み、サポート範囲を全域へ拡大。"
    ])
    
    # Card 2: Self-improving AI
    c2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), box_w, box_h)
    tf2 = c2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "🔄 使うほど賢くなるAI学習ループ"
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN
    p2.space_after = Pt(10)
    add_bullets(tf2, [
        "質問ログと「役に立ったか」の評価を記録。良質な回答を学習し、不要なものを除いて回答精度を高める仕組みを構築。"
    ])
    
    # Card 3: Executive Closing
    c3 = slide.shapes.add_textbox(Inches(0.75), Inches(4.3), Inches(11.833), Inches(2.0))
    tf3 = c3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "💡 まとめ"
    p3.font.name = FONT_TITLE
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = TEXT_LIGHT
    p3.space_after = Pt(10)
    add_bullets(tf3, [
        "このアプリは単なる点検表のデジタル化ではありません。現場がより安全に働き、故障を防ぎ、設備を最高の状態で稼働させ続けるためのツールです。"
    ])
    
    # Save presentation
    output_path = 'seibi_presentation.pptx'
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == '__main__':
    create_presentation()
